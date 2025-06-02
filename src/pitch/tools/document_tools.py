from crewai.tools import BaseTool
from typing import Type
from pydantic import BaseModel, Field
from pypdf import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import json
import os

class ParseDocumentInput(BaseModel):
    """Input schema for document parsing tool."""
    file_path: str = Field(..., description="Path to the document file to parse")

class WebResearchInput(BaseModel):
    """Input schema for web research tool."""
    query: str = Field(..., description="Search query about the startup or industry")

class DocumentParserTool(BaseTool):
    name: str = "Document Parser"
    description: str = (
        "Tool for parsing pitch deck documents (PDF/PPT) and extracting their content"
    )
    args_schema: Type[BaseModel] = ParseDocumentInput

    def _run(self, file_path: str) -> str:
        print(f"\nDocument Parser received file_path: {file_path}")
        
        # Convert to absolute path if not already
        abs_file_path = os.path.abspath(file_path)
        print(f"Absolute path: {abs_file_path}")
        
        # Verify file exists
        exists = os.path.exists(abs_file_path)
        print(f"File exists: {exists}")
        if not exists:
            raise ValueError(f"File not found: {abs_file_path}")
        
        # Check file extension
        is_pdf = abs_file_path.lower().endswith('.pdf')
        is_ppt = abs_file_path.lower().endswith(('.ppt', '.pptx'))
        print(f"Is PDF: {is_pdf}, Is PPT: {is_ppt}")
        
        if is_pdf:
            return self._parse_pdf(abs_file_path)
        elif is_ppt:
            return self._parse_ppt(abs_file_path)
        else:
            raise ValueError(f"Unsupported file format for {abs_file_path}. Only PDF and PPT/PPTX are supported.")

    def _parse_pdf(self, file_path: str) -> str:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _parse_ppt(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

class WebResearchTool(BaseTool):
    name: str = "Web Research"
    description: str = (
        "Tool for conducting web research about startups and industries"
    )
    args_schema: Type[BaseModel] = WebResearchInput

    def _run(self, query: str) -> str:
        try:
            # This is a simple example - in production you might want to use 
            # more sophisticated search APIs and handle rate limiting
            search_url = f"https://api.bing.microsoft.com/v7.0/search"
            headers = {
                "Ocp-Apim-Subscription-Key": "YOUR-KEY-HERE"  # Replace with actual key
            }
            params = {
                "q": query,
                "count": 5,
                "offset": 0,
                "mkt": "en-US"
            }
            
            response = requests.get(search_url, headers=headers, params=params)
            results = response.json()
            
            if "webPages" in results:
                pages = results["webPages"]["value"]
                research_results = []
                
                for page in pages:
                    try:
                        # Get the webpage content
                        page_response = requests.get(page["url"])
                        soup = BeautifulSoup(page_response.text, 'html.parser')
                        
                        # Extract main content (this is a simple example)
                        content = " ".join([p.get_text() for p in soup.find_all('p')])
                        
                        research_results.append({
                            "title": page["name"],
                            "url": page["url"],
                            "snippet": page["snippet"],
                            "content": content[:500]  # First 500 chars of content
                        })
                    except Exception as e:
                        continue
                        
                return json.dumps(research_results, indent=2)
            
            return "No results found"
            
        except Exception as e:
            return f"Error performing web research: {str(e)}"
